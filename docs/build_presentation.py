"""
Build MatRes hackathon pitch deck as a .pptx file.
Run: python docs/build_presentation.py
Output: docs/MatRes_Pitch_Deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ── Colours ────────────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT_BLUE  = RGBColor(0x42, 0x85, 0xF4)   # Google blue
ACCENT_GREEN = RGBColor(0x34, 0xA8, 0x53)   # Google green
ACCENT_RED   = RGBColor(0xEA, 0x43, 0x35)   # Google red
ACCENT_YELL  = RGBColor(0xFB, 0xBC, 0x04)   # Google yellow
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xCC, 0xCC, 0xCC)
MID_GREY     = RGBColor(0x88, 0x88, 0x99)
CARD_BG      = RGBColor(0x15, 0x28, 0x3E)   # slightly lighter navy for cards

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank


def bg(slide, color=DARK_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        bg_color=None, border_color=None, border_pt=0):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def accent_bar(slide, color=ACCENT_BLUE, left=0.4, top=1.1, width=1.5, height=0.06):
    """Thin coloured underline bar."""
    b = box(slide, left, top, width, height, bg_color=color)
    return b


def add_bullet_tf(slide, items, left, top, width, height,
                  size=16, color=LIGHT_GREY, bullet_color=ACCENT_BLUE,
                  bold_first=False):
    """Add a text frame with bullet points."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
    return txBox


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Top accent strip
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_BLUE)

# Google badge
txt(s, "Google for Startups AI Agents Challenge  ·  Track 1: Build  ·  2026",
    0.5, 0.18, 12, 0.35, size=11, color=MID_GREY)

# Main title
txt(s, "MatRes", 0.5, 1.1, 9, 1.4, size=80, bold=True, color=WHITE)
accent_bar(s, ACCENT_BLUE, 0.5, 2.55, 4.0, 0.07)

# Subtitle
txt(s, "Materials Resilience Agent", 0.5, 2.75, 9, 0.6, size=28, color=ACCENT_BLUE, bold=True)
txt(s, "Supply risk analysis for EV battery engineering teams",
    0.5, 3.4, 9, 0.5, size=20, color=LIGHT_GREY)
txt(s, "From BOM upload to full risk report — under 90 seconds.",
    0.5, 3.95, 9, 0.4, size=16, color=MID_GREY)

# Demo URL card
box(s, 0.5, 5.0, 8.5, 0.85, bg_color=CARD_BG)
txt(s, "🌐  Live demo:", 0.75, 5.1, 2.5, 0.5, size=13, color=MID_GREY)
txt(s, "https://matres-705351137331.us-central1.run.app",
    0.75, 5.38, 8.0, 0.4, size=14, bold=True, color=ACCENT_BLUE)

# Right side visual block
box(s, 10.0, 1.0, 3.0, 5.5, bg_color=CARD_BG)
txt(s, "⚗️", 10.9, 1.8, 1.5, 1.5, size=64, align=PP_ALIGN.CENTER)
txt(s, "AI · ADK · Gemini 2.5", 10.0, 3.5, 3.0, 0.4,
    size=11, color=MID_GREY, align=PP_ALIGN.CENTER)
txt(s, "Cloud Run · MCP", 10.0, 3.9, 3.0, 0.4,
    size=11, color=MID_GREY, align=PP_ALIGN.CENTER)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_RED)

txt(s, "THE PROBLEM", 0.5, 0.18, 4, 0.35, size=11, color=ACCENT_RED, bold=True)
txt(s, "One material disruption can cost $1.9 billion",
    0.5, 0.65, 11, 1.0, size=36, bold=True, color=WHITE)
accent_bar(s, ACCENT_RED, 0.5, 1.65, 5.0, 0.06)

incidents = [
    ("🔴  GM Bolt Recall 2021",
     "141,667 vehicles · $1.9B cost · Lithium battery thermal runaway"),
    ("🔴  China Graphite Export Ban — Dec 2023",
     "77% of global supply disrupted overnight · No advance warning"),
    ("🔴  DRC Cobalt Concentration",
     "73% of world supply · FEOC non-compliant under US IRA 2022"),
]
y = 1.9
for title, detail in incidents:
    box(s, 0.5, y, 11.5, 0.95, bg_color=CARD_BG)
    txt(s, title,  0.8, y + 0.08, 11, 0.38, size=15, bold=True, color=WHITE)
    txt(s, detail, 0.8, y + 0.46, 11, 0.38, size=13, color=LIGHT_GREY)
    y += 1.1

box(s, 0.5, 5.45, 11.5, 0.7, bg_color=RGBColor(0x2A, 0x10, 0x10))
txt(s, "⏱   Engineers spend 3–6 weeks doing this analysis manually. Every time.",
    0.8, 5.55, 11, 0.5, size=15, bold=True, color=ACCENT_YELL)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_RED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_GREEN)

txt(s, "THE SOLUTION", 0.5, 0.18, 4, 0.35, size=11, color=ACCENT_GREEN, bold=True)
txt(s, "Upload a BOM. Get a risk report in 90 seconds.",
    0.5, 0.65, 12, 0.8, size=34, bold=True, color=WHITE)
accent_bar(s, ACCENT_GREEN, 0.5, 1.5, 4.5, 0.06)

steps = [
    ("01", ACCENT_BLUE,  "Supply Risk",
     "Which materials are geopolitically concentrated\nHHI score + FEOC flag per component"),
    ("02", ACCENT_RED,   "Failure Modes",
     "What failures those materials caused in real vehicles\nReal NHTSA recalls, severity 1–5"),
    ("03", ACCENT_GREEN, "Substitutions",
     "3 ranked alternatives with full property delta\nEnergy density, cycle life, cost, CO₂"),
    ("04", ACCENT_YELL,  "Qualification Roadmap",
     "Standards, timeline in weeks, cost band in USD\nUN 38.3 · IEC 62660 · UL 2580 · ISO 26262"),
]
x = 0.4
for num, color, title, detail in steps:
    box(s, x, 1.75, 2.9, 3.8, bg_color=CARD_BG)
    box(s, x, 1.75, 2.9, 0.12, bg_color=color)
    txt(s, num,    x+0.1, 1.95, 0.6, 0.6, size=28, bold=True, color=color)
    txt(s, title,  x+0.1, 2.55, 2.7, 0.5, size=17, bold=True, color=WHITE)
    txt(s, detail, x+0.1, 3.1,  2.7, 2.0, size=12, color=LIGHT_GREY)
    x += 3.1

txt(s, "Every number in the output has a source citation. No uncited claims reach the UI.",
    0.5, 6.85, 12, 0.45, size=13, color=MID_GREY, align=PP_ALIGN.CENTER)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — GOOGLE TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_BLUE)

txt(s, "TECH STACK", 0.5, 0.18, 4, 0.35, size=11, color=ACCENT_BLUE, bold=True)
txt(s, "Built entirely on Google's AI stack",
    0.5, 0.65, 10, 0.8, size=34, bold=True, color=WHITE)
accent_bar(s, ACCENT_BLUE, 0.5, 1.5, 3.5, 0.06)

stack = [
    (ACCENT_BLUE,  "Gemini 2.5 Pro (Vertex AI)",
     "Executive summary generation + reasoning over multi-source BOM data.\nAll calls IAM-controlled — no direct API key in production."),
    (ACCENT_GREEN, "Google Agent Development Kit (ADK)",
     "Orchestrates 4 sub-agents with typed inputs, outputs, and FunctionTool calls.\nRoot agent aggregates outputs into a single RiskReport."),
    (ACCENT_YELL,  "Model Context Protocol (MCP)",
     "USGS MCS 2025 and Materials Project data wrapped as standalone MCP servers.\nAgent queries them via tool calls — no hardcoded data access."),
    (ACCENT_RED,   "Cloud Run + Secret Manager + Cloud Build",
     "Fully serverless, GCP us-central1. API keys in Secret Manager — zero\nhardcoded credentials in the container image or git history."),
]

y = 1.75
for color, title, detail in stack:
    box(s, 0.5, y, 0.12, 0.75, bg_color=color)
    txt(s, title,  0.78, y,        11.5, 0.38, size=15, bold=True, color=WHITE)
    txt(s, detail, 0.78, y + 0.38, 11.5, 0.42, size=12, color=LIGHT_GREY)
    y += 1.1

box(s, 0.5, 6.3, 11.5, 0.55, bg_color=CARD_BG)
txt(s, "💡  Gemini 2.5 Pro routes through Vertex AI — auditable, billed to GCP project, IAM-controlled",
    0.75, 6.38, 11, 0.4, size=12, color=ACCENT_YELL)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_BLUE)

txt(s, "ARCHITECTURE", 0.5, 0.18, 4, 0.35, size=11, color=ACCENT_BLUE, bold=True)
txt(s, "4 ADK sub-agents · MCP servers · Hallucination guard · Cloud Run",
    0.5, 0.6, 12, 0.5, size=18, color=LIGHT_GREY)

arch_png = os.path.join(os.path.dirname(__file__), "architecture.png")
if os.path.exists(arch_png):
    s.shapes.add_picture(arch_png, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))
else:
    box(s, 0.5, 1.2, 12.3, 5.9, bg_color=CARD_BG)
    txt(s, "[architecture.png not found — place docs/architecture.png]",
        3.0, 3.5, 7, 0.5, size=14, color=MID_GREY, align=PP_ALIGN.CENTER)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DEMO RESULT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_YELL)

txt(s, "LIVE DEMO — SCENARIO A", 0.5, 0.18, 5, 0.35, size=11, color=ACCENT_YELL, bold=True)
txt(s, "NMC 811 Battery Pack — what MatRes finds in 14 seconds",
    0.5, 0.65, 12, 0.7, size=30, bold=True, color=WHITE)
accent_bar(s, ACCENT_YELL, 0.5, 1.4, 5.0, 0.06)

# Left — risks
box(s, 0.4, 1.6, 5.9, 5.5, bg_color=CARD_BG)
box(s, 0.4, 1.6, 5.9, 0.38, bg_color=ACCENT_RED)
txt(s, "SUPPLY RISK IDENTIFIED", 0.6, 1.66, 5.5, 0.3, size=12, bold=True, color=WHITE)

risks = [
    ("Cobalt",   "HIGH",   "HHI 5,535 · DRC 73% · FEOC ⚠️"),
    ("Graphite", "HIGH",   "HHI 6,058 · China 77% · Export ban active"),
    ("Nickel",   "HIGH",   "HHI 3,576 · Indonesia 55%"),
    ("Lithium",  "HIGH",   "HHI 3,142 · Australia 47%"),
]
y = 2.15
for mat, level, detail in risks:
    txt(s, f"🔴  {mat}", 0.65, y, 2.5, 0.32, size=14, bold=True, color=WHITE)
    txt(s, detail,        0.65, y+0.32, 5.4, 0.28, size=11, color=LIGHT_GREY)
    y += 0.72

# Right — substitution
box(s, 6.7, 1.6, 6.2, 5.5, bg_color=CARD_BG)
box(s, 6.7, 1.6, 6.2, 0.38, bg_color=ACCENT_GREEN)
txt(s, "TOP SUBSTITUTION RECOMMENDED", 6.9, 1.66, 5.8, 0.3, size=12, bold=True, color=WHITE)

txt(s, "#1  LFP (LiFePO₄)", 6.9, 2.15, 5.8, 0.5, size=22, bold=True, color=ACCENT_GREEN)
txt(s, "Composite score: 67 / 100", 6.9, 2.65, 5.8, 0.35, size=14, color=LIGHT_GREY)

metrics = [
    ("Supply risk",    "10 / 100",       "(vs 70+ for cobalt)", ACCENT_GREEN),
    ("Energy delta",   "−42%",           "trade-off vs NMC 811", ACCENT_YELL),
    ("Qualification",  "80 weeks",       "~19 months",           LIGHT_GREY),
    ("Cost band",      "$925k–$1.85M",   "full certification",   LIGHT_GREY),
]
y = 3.2
for label, value, note, color in metrics:
    txt(s, label,  6.9,  y,      2.5, 0.32, size=12, color=MID_GREY)
    txt(s, value,  9.2,  y,      2.2, 0.32, size=14, bold=True, color=color)
    txt(s, note,   11.2, y,      1.6, 0.32, size=11, color=MID_GREY)
    y += 0.55

txt(s, "Source: USGS MCS 2025 · NREL Battery Report 2023 · IEC 62660",
    6.9, 6.7, 5.8, 0.3, size=10, color=MID_GREY)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_YELL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — HALLUCINATION GUARD
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_RED)

txt(s, "HALLUCINATION GUARD", 0.5, 0.18, 5, 0.35, size=11, color=ACCENT_RED, bold=True)
txt(s, "Every number is cited. By design.",
    0.5, 0.65, 10, 0.7, size=34, bold=True, color=WHITE)
accent_bar(s, ACCENT_RED, 0.5, 1.4, 3.5, 0.06)

txt(s, "Engineers cannot act on uncited AI output — legal and regulatory exposure.",
    0.5, 1.6, 12, 0.4, size=15, color=LIGHT_GREY)

# Table header
box(s, 0.4, 2.2, 5.8, 0.45, bg_color=RGBColor(0x3A, 0x10, 0x10))
box(s, 6.5, 2.2, 6.4, 0.45, bg_color=RGBColor(0x0D, 0x30, 0x1A))
txt(s, "❌  Without MatRes", 0.6, 2.3, 5.5, 0.3, size=13, bold=True, color=ACCENT_RED)
txt(s, "✅  With MatRes",    6.7, 2.3, 6.0, 0.3, size=13, bold=True, color=ACCENT_GREEN)

rows = [
    ('"Cobalt is risky"',
     '"Cobalt HHI: 5,535  ·  Source: USGS MCS 2025"'),
    ('"LFP is safer"',
     '"LFP supply risk: 10/100  ·  Source: NREL 2023"'),
    ('"Qualification takes months"',
     '"80 weeks · $925k–$1.85M  ·  Source: UN 38.3, IEC 62660"'),
    ('"China dominates graphite"',
     '"China: 77% share · HHI 6,058  ·  Source: USGS MCS 2025"'),
]
y = 2.85
for bad, good in rows:
    box(s, 0.4,  y, 5.8, 0.7, bg_color=CARD_BG)
    box(s, 6.5,  y, 6.4, 0.7, bg_color=CARD_BG)
    box(s, 6.25, y, 0.25, 0.7, bg_color=RGBColor(0x15, 0x28, 0x3E))
    txt(s, bad,  0.6, y+0.12, 5.5, 0.5, size=12, color=MID_GREY)
    txt(s, good, 6.65, y+0.1, 6.1, 0.5, size=12, color=ACCENT_GREEN)
    y += 0.82

box(s, 0.4, 6.3, 12.5, 0.55, bg_color=RGBColor(0x20, 0x10, 0x10))
txt(s, "CitationError middleware rejects any numeric output without a source field — before it reaches the UI",
    0.65, 6.38, 12, 0.4, size=12, color=ACCENT_YELL)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_RED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WHY NET-NEW AGENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_BLUE)

txt(s, "TRACK 1: BUILD", 0.5, 0.18, 4, 0.35, size=11, color=ACCENT_BLUE, bold=True)
txt(s, "A genuinely new agent — not a wrapper",
    0.5, 0.65, 10, 0.7, size=34, bold=True, color=WHITE)
accent_bar(s, ACCENT_BLUE, 0.5, 1.4, 4.0, 0.06)

# Left: NOT column
box(s, 0.4, 1.65, 5.9, 4.8, bg_color=CARD_BG)
box(s, 0.4, 1.65, 5.9, 0.38, bg_color=RGBColor(0x3A, 0x10, 0x10))
txt(s, "What it is NOT", 0.6, 1.71, 5.5, 0.28, size=12, bold=True, color=ACCENT_RED)
nots = [
    "Not a supply chain dashboard\n(Interos, Resilinc operate at KPI level)",
    "Not a data lookup tool\n(USGS, Materials Project have no agent layer)",
    "Not a general LLM query\n(ChatGPT/Gemini give uncited estimates)",
]
y = 2.2
for item in nots:
    txt(s, "✗  " + item, 0.65, y, 5.4, 0.85, size=13, color=LIGHT_GREY)
    y += 1.1

# Right: IS column
box(s, 6.7, 1.65, 6.2, 4.8, bg_color=CARD_BG)
box(s, 6.7, 1.65, 6.2, 0.38, bg_color=RGBColor(0x0D, 0x30, 0x1A))
txt(s, "What makes it novel", 6.9, 1.71, 5.8, 0.28, size=12, bold=True, color=ACCENT_GREEN)
iss = [
    "Connects geopolitical risk → failure data\n→ substitution → qualification in one pipeline",
    "Hallucination guard enforces citations\nfor regulatory-grade engineering output",
    "Deterministic sub-agents (FunctionTool)\nFast, auditable, no hallucination on data",
    "Composite scorer with tunable weights\nEngineers adjust cost vs performance vs CO₂",
]
y = 2.2
for item in iss:
    txt(s, "✓  " + item, 6.9, y, 5.8, 0.85, size=13, color=LIGHT_GREY)
    y += 1.1

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BUSINESS CASE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_GREEN)

txt(s, "BUSINESS CASE", 0.5, 0.18, 4, 0.35, size=11, color=ACCENT_GREEN, bold=True)
txt(s, "Replacing 3 weeks of manual work with 90 seconds",
    0.5, 0.65, 11, 0.7, size=32, bold=True, color=WHITE)
accent_bar(s, ACCENT_GREEN, 0.5, 1.4, 5.0, 0.06)

cards = [
    (ACCENT_BLUE,  "TARGET CUSTOMERS",
     "Tier-1 EV battery suppliers\nPanasonic · LG Energy · Samsung SDI · CATL\n\nOEM procurement teams\nGM · Ford · Tata · Ola Electric\n\nRegulatory compliance teams\nUS IRA 2022 FEOC deadline reviews"),
    (ACCENT_GREEN, "MARKET OPPORTUNITY",
     "$4B supply chain risk software market\n\nNo agentic competitor at the\nengineering-decision layer today\n\nExisting tools operate at dashboard level —\nnot material-decision specific"),
    (ACCENT_YELL,  "POST-HACKATHON PATH",
     "Expand beyond EV batteries:\n· Semiconductor BOMs\n· Pharmaceutical supply chains\n· Aerospace materials\n\nEnterprise SaaS or\nAPI integration into ERP/PLM systems"),
]
x = 0.4
for color, title, content in cards:
    box(s, x, 1.65, 3.95, 5.5, bg_color=CARD_BG)
    box(s, x, 1.65, 3.95, 0.38, bg_color=color)
    txt(s, title,   x+0.15, 1.71, 3.65, 0.28, size=12, bold=True, color=WHITE)
    txt(s, content, x+0.15, 2.15, 3.65, 4.8,  size=13, color=LIGHT_GREY)
    x += 4.3

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.08, bg_color=ACCENT_BLUE)

txt(s, "TRY IT NOW", 0.5, 0.18, 4, 0.35, size=11, color=ACCENT_BLUE, bold=True)
txt(s, "Live now — test it yourself in 2 minutes",
    0.5, 0.65, 10, 0.7, size=32, bold=True, color=WHITE)
accent_bar(s, ACCENT_BLUE, 0.5, 1.4, 4.5, 0.06)

# Demo URL — large
box(s, 0.4, 1.65, 12.5, 0.85, bg_color=CARD_BG)
txt(s, "🌐   https://matres-705351137331.us-central1.run.app",
    0.65, 1.75, 12, 0.55, size=20, bold=True, color=ACCENT_BLUE)

# Steps
steps_data = [
    ("01", "Download test BOM",
     "Get bom_fixtures/nmc811.json from github.com/raunakmantri9/MatRes"),
    ("02", "Upload in the sidebar",
     "Open the demo URL → use the sidebar file uploader"),
    ("03", "See the full report",
     "Cobalt flagged HIGH → LFP ranked #1 → 19-month qualification roadmap"),
]
x = 0.4
for num, title, detail in steps_data:
    box(s, x, 2.75, 3.95, 2.5, bg_color=CARD_BG)
    txt(s, num,    x+0.15, 2.88, 1.0,  0.55, size=32, bold=True, color=ACCENT_BLUE)
    txt(s, title,  x+0.15, 3.48, 3.65, 0.45, size=15, bold=True, color=WHITE)
    txt(s, detail, x+0.15, 3.95, 3.65, 1.1,  size=12, color=LIGHT_GREY)
    x += 4.3

# GitHub
box(s, 0.4, 5.5, 12.5, 0.55, bg_color=CARD_BG)
txt(s, "⌨️   github.com/raunakmantri9/MatRes",
    0.65, 5.58, 12, 0.38, size=14, color=LIGHT_GREY)

# Footer
txt(s, "Built by Raunak Mantri  ·  Google for Startups AI Agents Challenge 2026  ·  Gemini 2.5 Pro · Google ADK · MCP · Cloud Run",
    0.4, 6.4, 12.5, 0.4, size=11, color=MID_GREY, align=PP_ALIGN.CENTER)

box(s, 0, 7.42, 13.33, 0.08, bg_color=ACCENT_BLUE)


# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "MatRes_Pitch_Deck.pptx")
prs.save(out)
print(f"Saved: {out}")
